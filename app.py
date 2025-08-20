# app_sih.py
import streamlit as st
import fitz  # PyMuPDF
import os
from src.keyword_extractor import extract_keywords_hierarchical

# Additional imports for DOCX and PPTX
from docx import Document
from pptx import Presentation

# -------------------------
# Streamlit App Layout
# -------------------------
st.set_page_config(page_title="AutoTagLearn – SIH Demo", layout="wide")

st.title("📚 AutoTagLearn ")
st.markdown("""
**AutoTagLearn: An NLP-Based System for Intelligent Keyword Extraction in Academic Content.**  

---
**Steps:**  
1. Upload a document (PDF, DOCX, PPTX) or video  
2. View extracted keywords  
3. Download keywords
""")

# -------------------------
# Upload Section
# -------------------------
uploaded_file = st.file_uploader(
    "📄 Upload a PDF, DOCX, PPTX or a video (mp4/mov/avi)", 
    type=["pdf", "docx", "pptx", "mp4", "mov", "avi"]
)

# -------------------------
# Helper: Flatten keywords
# -------------------------
def flatten_keywords(keywords_dict):
    flat = []
    for main_kw, sub_kws in keywords_dict.items():
        for kw in sub_kws:
            if kw not in flat:
                flat.append(kw)
    return flat

# -------------------------
# Helper: Display hierarchy
# -------------------------
def display_keywords_hierarchy(keywords_dict, section_name="Keywords"):
    st.subheader(f"🔑 {section_name}")
    if not keywords_dict:
        st.info("No keywords extracted.")
        return
    for main_kw, sub_kws in keywords_dict.items():
        st.markdown(f"- **{main_kw}**")
        for sub_kw in sub_kws:
            if sub_kw != main_kw:
                st.markdown(f" • {sub_kw}")

# -------------------------
# Extract text from DOCX
# -------------------------
def extract_docx_text(path):
    doc = Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

# -------------------------
# Extract text from PPTX
# -------------------------
def extract_pptx_text(path):
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

# -------------------------
# Processing Section
# -------------------------
if uploaded_file:
    ext = uploaded_file.name.split(".")[-1].lower()
    os.makedirs("data/uploads", exist_ok=True)
    file_path = os.path.join("data/uploads", uploaded_file.name)
    with open(file_path, "wb") as f:
        f.write(uploaded_file.read())

    if ext == "pdf":
        doc = fitz.open(file_path)
        text = "".join([page.get_text() for page in doc])
        section_name = "PDF Keywords"

    elif ext == "docx":
        text = extract_docx_text(file_path)
        section_name = "DOCX Keywords"

    elif ext == "pptx":
        text = extract_pptx_text(file_path)
        section_name = "PPTX Keywords"

    elif ext in ["mp4", "mov", "avi"]:
        import whisper
        model = whisper.load_model("tiny")
        result = model.transcribe(file_path)
        text = result['text']
        section_name = "Video Keywords"

    else:
        st.error("Unsupported file type!")
        st.stop()

    st.subheader(f"📄 Extracted Text (first 300 chars)")
    st.write(text[:300] + "...")

    keywords = extract_keywords_hierarchical(text, top_n=10)
    display_keywords_hierarchy(keywords, section_name)

    flat_keywords = flatten_keywords(keywords)
    st.download_button(
        "⬇ Download Keywords",
        "\n".join(flat_keywords),
        file_name=f"{uploaded_file.name}_keywords.txt"
    )
